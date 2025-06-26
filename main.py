import os
import tempfile
from typing import List, Dict, Any, Optional
from pathlib import Path

from fastapi import FastAPI, File, UploadFile, Form, HTTPException
from fastapi.responses import JSONResponse, FileResponse
from fastapi.middleware.cors import CORSMiddleware
import whisper
import openai
from pptx import Presentation
from pptx.util import Inches
from dotenv import load_dotenv

load_dotenv()

app = FastAPI(title="SlideSynth API", version="1.0.0")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

whisper_model = whisper.load_model("base")
openai.api_key = os.getenv("OPENAI_API_KEY")

class SlideGenerator:
    @staticmethod
    def generate_slides(content: str) -> List[Dict[str, Any]]:
        prompt = f"""
        Transform the following content into a structured 10-slide investor pitch deck. 
        Return ONLY a JSON array with exactly 10 slides, each containing:
        - title: slide title (string)
        - bullets: array of 3-5 bullet points (strings)
        
        Content: {content}
        
        The 10 slides should follow this structure:
        1. Problem Statement
        2. Solution Overview
        3. Market Opportunity
        4. Product/Service Details
        5. Business Model
        6. Competitive Analysis
        7. Marketing Strategy
        8. Financial Projections
        9. Team & Execution
        10. Funding Ask & Use of Funds
        """
        
        try:
            client = openai.OpenAI(api_key=os.getenv("OPENAI_API_KEY"))
            response = client.chat.completions.create(
                model="gpt-4",
                messages=[{"role": "user", "content": prompt}],
                temperature=0.7
            )
            
            import json
            slides_data = json.loads(response.choices[0].message.content)
            return slides_data
        except Exception as e:
            print(f"Error generating slides with OpenAI: {e}")
            print(f"Content length: {len(content)}")
            print(f"Content preview: {content[:100]}...")
            return SlideGenerator._get_dummy_slides()
    
    @staticmethod
    def _get_dummy_slides() -> List[Dict[str, Any]]:
        return [
            {
                "title": "Problem Statement",
                "bullets": [
                    "Current market pain point identification",
                    "Quantified impact on target users",
                    "Existing solution limitations",
                    "Market gap opportunity"
                ]
            },
            {
                "title": "Solution Overview",
                "bullets": [
                    "Core product functionality",
                    "Key differentiating features",
                    "User experience benefits",
                    "Technology approach"
                ]
            },
            {
                "title": "Market Opportunity",
                "bullets": [
                    "Total addressable market size",
                    "Target customer segments",
                    "Market growth trends",
                    "Revenue potential"
                ]
            },
            {
                "title": "Product Details",
                "bullets": [
                    "Feature specifications",
                    "Development roadmap",
                    "Technical architecture",
                    "User interface design"
                ]
            },
            {
                "title": "Business Model",
                "bullets": [
                    "Revenue streams",
                    "Pricing strategy",
                    "Customer acquisition cost",
                    "Unit economics"
                ]
            },
            {
                "title": "Competitive Analysis",
                "bullets": [
                    "Direct competitors",
                    "Indirect competitors",
                    "Competitive advantages",
                    "Market positioning"
                ]
            },
            {
                "title": "Marketing Strategy",
                "bullets": [
                    "Go-to-market approach",
                    "Customer acquisition channels",
                    "Brand positioning",
                    "Partnership opportunities"
                ]
            },
            {
                "title": "Financial Projections",
                "bullets": [
                    "3-year revenue forecast",
                    "Key financial metrics",
                    "Break-even analysis",
                    "Growth assumptions"
                ]
            },
            {
                "title": "Team & Execution",
                "bullets": [
                    "Founding team background",
                    "Key advisors",
                    "Hiring roadmap",
                    "Execution milestones"
                ]
            },
            {
                "title": "Funding Ask",
                "bullets": [
                    "Investment amount sought",
                    "Use of funds breakdown",
                    "Expected milestones",
                    "Exit strategy timeline"
                ]
            }
        ]

class PPTXGenerator:
    @staticmethod
    def create_presentation(slides_data: List[Dict[str, Any]], filename: str) -> str:
        prs = Presentation()
        
        for slide_data in slides_data:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            
            title = slide.shapes.title
            title.text = slide_data["title"]
            
            content = slide.placeholders[1]
            text_frame = content.text_frame
            text_frame.clear()
            
            for bullet in slide_data["bullets"]:
                p = text_frame.add_paragraph()
                p.text = bullet
                p.level = 0
        
        filepath = f"/tmp/{filename}"
        prs.save(filepath)
        return filepath

@app.get("/")
async def root():
    return {"message": "SlideSynth API"}

@app.post("/transcribe")
async def transcribe_audio(file: UploadFile = File(...)):
    if not file.content_type.startswith("audio/"):
        raise HTTPException(status_code=400, detail="File must be audio format")
    
    with tempfile.NamedTemporaryFile(delete=False, suffix=".wav") as tmp_file:
        content = await file.read()
        tmp_file.write(content)
        tmp_file.flush()
        
        try:
            result = whisper_model.transcribe(tmp_file.name)
            transcription = result["text"]
        finally:
            os.unlink(tmp_file.name)
    
    return {"transcription": transcription}

@app.post("/generate-slides")
async def generate_slides(
    text_content: Optional[str] = Form(None),
    audio_file: Optional[UploadFile] = File(None)
):
    if not text_content and not audio_file:
        raise HTTPException(status_code=400, detail="Either text_content or audio_file must be provided")
    
    content = text_content
    
    if audio_file:
        if not audio_file.content_type.startswith("audio/"):
            raise HTTPException(status_code=400, detail="File must be audio format")
        
        with tempfile.NamedTemporaryFile(delete=False, suffix=".wav") as tmp_file:
            audio_content = await audio_file.read()
            tmp_file.write(audio_content)
            tmp_file.flush()
            
            try:
                result = whisper_model.transcribe(tmp_file.name)
                content = result["text"]
            finally:
                os.unlink(tmp_file.name)
    
    slides = SlideGenerator.generate_slides(content)
    
    return {
        "slides": slides,
        "slide_count": len(slides),
        "source_content": content[:200] + "..." if len(content) > 200 else content
    }

@app.post("/generate-pptx")
async def generate_pptx(
    text_content: Optional[str] = Form(None),
    audio_file: Optional[UploadFile] = File(None),
    filename: str = Form("slidesynth_presentation.pptx")
):
    if not text_content and not audio_file:
        raise HTTPException(status_code=400, detail="Either text_content or audio_file must be provided")
    
    content = text_content
    
    if audio_file:
        if not audio_file.content_type.startswith("audio/"):
            raise HTTPException(status_code=400, detail="File must be audio format")
        
        with tempfile.NamedTemporaryFile(delete=False, suffix=".wav") as tmp_file:
            audio_content = await audio_file.read()
            tmp_file.write(audio_content)
            tmp_file.flush()
            
            try:
                result = whisper_model.transcribe(tmp_file.name)
                content = result["text"]
            finally:
                os.unlink(tmp_file.name)
    
    slides = SlideGenerator.generate_slides(content)
    pptx_path = PPTXGenerator.create_presentation(slides, filename)
    
    return FileResponse(
        path=pptx_path,
        filename=filename,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)